"""ゼミ進捗報告スライド生成スクリプト."""

from __future__ import annotations

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# ── テーマカラー ──
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)
BG_MID = RGBColor(0x22, 0x22, 0x3A)
ACCENT = RGBColor(0x4E, 0xC9, 0xB0)
ACCENT2 = RGBColor(0x56, 0x9C, 0xD6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
ORANGE = RGBColor(0xFF, 0xA5, 0x00)
RED_SOFT = RGBColor(0xE0, 0x60, 0x60)
GREEN_SOFT = RGBColor(0x60, 0xE0, 0x80)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name="Meiryo"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return tf


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=LIGHT_GRAY, bullet_color=ACCENT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Meiryo"
        p.space_after = Pt(6)
        p.level = 0
    return tf


def add_rounded_rect(slide, left, top, width, height, fill_color, text="",
                     font_size=12, font_color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.name = "Meiryo"
        p.font.bold = True
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    shape.text_frame.auto_size = None
    shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape


def add_arrow(slide, x1, y1, x2, y2, color=ACCENT):
    connector = slide.shapes.add_connector(
        1, x1, y1, x2, y2  # 1 = straight connector
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(2)


def slide_title(slide, title, subtitle=""):
    set_bg(slide, BG_DARK)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
                 title, font_size=32, color=WHITE, bold=True)
    # accent line
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.8), Inches(1.05), Inches(2), Pt(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
                     subtitle, font_size=14, color=LIGHT_GRAY)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]  # blank layout

    # ════════════════════════════════════════
    # Slide 1: 表紙
    # ════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    set_bg(s, BG_DARK)
    add_text_box(s, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
                 "LLMによるドローン制御コードの\n自動改変・評価基盤の構築",
                 font_size=36, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # accent line
    shape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(4.5), Inches(3.3), Inches(4), Pt(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    add_text_box(s, Inches(1), Inches(3.6), Inches(11), Inches(0.5),
                 "前田研究室 ゼミ進捗報告", font_size=20, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(1), Inches(4.3), Inches(11), Inches(0.5),
                 "2026年3月17日", font_size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    # ════════════════════════════════════════
    # Slide 2: 目次
    # ════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    slide_title(s, "目次")
    items = [
        "1. 研究概要",
        "2. 現在の進捗状況",
        "3. システムアーキテクチャ",
        "4. コードベース構成と各層の役割",
        "5. End-to-End 動作確認結果",
        "6. Kubernetes案の検討と却下",
        "7. 評価方法",
        "8. ネクストアクション",
    ]
    add_bullet_list(s, Inches(1.5), Inches(1.8), Inches(9), Inches(5),
                    items, font_size=22, color=WHITE)

    # ════════════════════════════════════════
    # Slide 3: 研究概要
    # ════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    slide_title(s, "研究概要")
    add_text_box(s, Inches(0.8), Inches(1.6), Inches(11), Inches(0.8),
                 "LLM（Claude Code等）を用いてドローン制御コードを自動生成・改変し、\n"
                 "CoppeliaSim上で評価して改善ループを回す基盤を構築する。",
                 font_size=18, color=LIGHT_GRAY)

    # 3つのボックス
    labels = [
        ("Controller\n改変", "Claude Code / Skill で\ncontroller/ を修正"),
        ("Simulation\n評価", "CoppeliaSim 上で\nbaseline vs candidate を比較"),
        ("結果分析\nフィードバック", "metrics / summary から\n次の改善方針を決定"),
    ]
    for i, (title, desc) in enumerate(labels):
        x = Inches(1.0 + i * 4.0)
        add_rounded_rect(s, x, Inches(3.0), Inches(3.2), Inches(1.0),
                         ACCENT2, title, font_size=16)
        add_text_box(s, x, Inches(4.2), Inches(3.2), Inches(1.0),
                     desc, font_size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    # ════════════════════════════════════════
    # Slide 4: 現在の進捗状況
    # ════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    slide_title(s, "現在の進捗状況")

    done_items = [
        "Agent Runner（run オーケストレーション）実装完了",
        "sim-eval（baseline/candidate 評価）実装完了",
        "Lock機構・Artifact Storage 実装完了",
        "Patch Provider（手動patch / API fallback）実装完了",
        "評価メトリクス算出・比較ロジック実装完了",
        "Agent Skills（eval-run / patch-generate / results-analyze）作成完了",
        "ユニットテスト 15件 全パス、lint パス",
    ]
    # Done column
    add_text_box(s, Inches(0.8), Inches(1.5), Inches(3), Inches(0.5),
                 "完了", font_size=20, color=GREEN_SOFT, bold=True)
    add_bullet_list(s, Inches(0.8), Inches(2.1), Inches(5.5), Inches(4.5),
                    done_items, font_size=14, color=LIGHT_GRAY)

    # Remaining column
    remaining = [
        "CoppeliaSim 実機接続テスト",
        "eval/run.py のスタブ→実接続への差し替え",
        "End-to-end の1 run 完走確認",
    ]
    add_text_box(s, Inches(7.5), Inches(1.5), Inches(3), Inches(0.5),
                 "未完了", font_size=20, color=ORANGE, bold=True)
    add_bullet_list(s, Inches(7.5), Inches(2.1), Inches(5), Inches(4),
                    remaining, font_size=14, color=LIGHT_GRAY)

    # ════════════════════════════════════════
    # Slide 5: アーキテクチャ図（現行版）
    # ════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    slide_title(s, "システムアーキテクチャ（現行版）", "Kubernetes なし・単一ホスト構成")

    # Operator box
    add_rounded_rect(s, Inches(0.5), Inches(1.8), Inches(2.5), Inches(1.0),
                     RGBColor(0x6A, 0x5A, 0xCD), "Operator\n(Claude Code / Skill)", font_size=12)

    # Agent Runner box
    add_rounded_rect(s, Inches(3.8), Inches(1.8), Inches(3.0), Inches(1.5),
                     RGBColor(0x2D, 0x5F, 0x8A),
                     "Agent Runner\n- lock / clone\n- patch apply\n"
                     "- static check\n- artifact collect",
                     font_size=11)

    # sim-eval box
    add_rounded_rect(s, Inches(7.8), Inches(1.8), Inches(2.8), Inches(1.5),
                     RGBColor(0x2D, 0x5F, 0x8A),
                     "sim-eval\n- baseline eval\n- candidate eval\n- metrics / summary",
                     font_size=11)

    # CoppeliaSim box
    add_rounded_rect(s, Inches(8.0), Inches(4.2), Inches(2.5), Inches(1.0),
                     RGBColor(0x8B, 0x45, 0x13), "CoppeliaSim\n(ZeroMQ / stepping)", font_size=12)

    # Local FS box
    add_rounded_rect(s, Inches(3.8), Inches(4.2), Inches(3.0), Inches(1.2),
                     RGBColor(0x3A, 0x3A, 0x5C),
                     "Local Filesystem\n/artifacts/runs/<run_id>/\n/workspace/<run_id>/",
                     font_size=11)

    # Arrows (using text-based arrows since connectors are limited)
    add_text_box(s, Inches(3.0), Inches(2.0), Inches(1.0), Inches(0.5),
                 "-->", font_size=20, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(6.8), Inches(2.0), Inches(1.0), Inches(0.5),
                 "-->", font_size=20, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(9.0), Inches(3.3), Inches(0.5), Inches(0.8),
                 "|", font_size=24, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(5.0), Inches(3.3), Inches(0.5), Inches(0.8),
                 "|", font_size=24, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

    # Label
    add_text_box(s, Inches(9.5), Inches(3.5), Inches(2), Inches(0.4),
                 "TCP/ZeroMQ", font_size=11, color=ACCENT)

    # Execution Host border label
    add_text_box(s, Inches(3.5), Inches(5.6), Inches(4), Inches(0.4),
                 "Execution Host（単一ホスト）", font_size=13, color=LIGHT_GRAY,
                 align=PP_ALIGN.CENTER)

    # ════════════════════════════════════════
    # Slide 6: コードベース構成と各層の役割
    # ════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    slide_title(s, "コードベース構成と各層の役割")

    # ── 左側: ディレクトリツリー ──
    tree_text = (
        "maeda_-laboratory/\n"
        "  agent_runner/\n"
        "    cli.py\n"
        "    runner.py\n"
        "    config.py\n"
        "    lock.py\n"
        "    artifacts.py\n"
        "    workspace.py\n"
        "    patch_provider.py\n"
        "    patch.py\n"
        "    static_check.py\n"
        "    models.py\n"
        "    llm.py\n"
        "    run_id.py\n"
        "  controller/\n"
        "    drone_controller.py\n"
        "  eval/\n"
        "    run.py\n"
        "    scenes.yaml\n"
        "  sim_eval/\n"
        "    cli.py\n"
        "    sim_client.py\n"
        "    evaluator.py\n"
        "    metrics.py\n"
        "    comparison.py\n"
        "  tests/\n"
        "  docs/"
    )
    add_text_box(s, Inches(0.5), Inches(1.5), Inches(3.8), Inches(5.8),
                 tree_text, font_size=11, color=ACCENT, font_name="Courier New")

    # ── 右側: 各層の説明カード ──
    layers = [
        (
            "agent_runner/",
            RGBColor(0x2D, 0x5F, 0x8A),
            "Run オーケストレーション層",
            "CLI エントリ / lock 排他制御 / fresh clone\n"
            "patch 受領・適用 / static check (lint,type,unit)\n"
            "sim-eval 起動 / artifact 収集・保存",
        ),
        (
            "controller/",
            RGBColor(0x6A, 0x5A, 0xCD),
            "LLM 改変対象（ドメイン層）",
            "DroneController クラス（制御ロジック本体）\n"
            "LLM / Operator が修正し patch.diff を生成する対象\n"
            "eval/ とは独立して変更可能",
        ),
        (
            "eval/",
            RGBColor(0x8B, 0x6B, 0x13),
            "評価実行契約（repo 契約）",
            "python -m eval.run で呼び出される評価エントリ\n"
            "scenes.yaml で scene_id → scene_path をマッピング\n"
            "PoC では固定（改変しない）",
        ),
        (
            "sim_eval/",
            RGBColor(0x2D, 0x7F, 0x5F),
            "シミュレーション評価層",
            "CoppeliaSim ZeroMQ 接続 (sim_client.py)\n"
            "baseline/candidate episode 実行 (evaluator.py)\n"
            "metrics 算出 → comparison → summary 生成",
        ),
        (
            "tests/ + docs/",
            RGBColor(0x5A, 0x5A, 0x6A),
            "品質保証 + 設計文書",
            "ユニットテスト 15件 (lock, metrics, comparison, patch)\n"
            "要件定義書 (not-k8s-architecture.md)",
        ),
    ]

    card_x = Inches(4.6)
    card_w = Inches(8.2)
    start_y = 1.45
    card_h_list = [1.1, 1.0, 1.0, 1.1, 0.8]

    for i, (dir_name, color, subtitle, desc) in enumerate(layers):
        y = Inches(start_y + sum(card_h_list[:i]) + i * 0.08)
        h = Inches(card_h_list[i])

        # カード背景
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  card_x, y, card_w, h)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0x28, 0x28, 0x40)
        card.line.color.rgb = color
        card.line.width = Pt(2)

        # ディレクトリ名ラベル
        label = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(4.75), y + Pt(6),
                                   Inches(1.8), Inches(0.32))
        label.fill.solid()
        label.fill.fore_color.rgb = color
        label.line.fill.background()
        tf = label.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = dir_name
        p.font.size = Pt(11)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.font.name = "Courier New"
        p.alignment = PP_ALIGN.CENTER

        # サブタイトル
        add_text_box(s, Inches(6.7), y + Pt(4),
                     Inches(5.8), Inches(0.3),
                     subtitle, font_size=13, color=WHITE, bold=True)

        # 説明文
        add_text_box(s, Inches(4.85), y + Inches(0.38),
                     Inches(7.7), h - Inches(0.4),
                     desc, font_size=10, color=LIGHT_GRAY)

    # ════════════════════════════════════════
    # Slide 7: End-to-End 動作確認結果
    # ════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    slide_title(s, "End-to-End 動作確認結果", "スタブ sim-eval での全パイプライン完走を確認")

    # パイプラインフロー（横並びボックス）
    pipeline_stages = [
        ("Clone", RGBColor(0x3A, 0x6A, 0x8A)),
        ("Patch\nApply", RGBColor(0x6A, 0x5A, 0xCD)),
        ("Static\nCheck", RGBColor(0x2D, 0x7F, 0x5F)),
        ("sim-eval", RGBColor(0x8B, 0x6B, 0x13)),
        ("Artifacts\n保存", RGBColor(0x2D, 0x5F, 0x8A)),
    ]
    for i, (label, color) in enumerate(pipeline_stages):
        x = Inches(0.8 + i * 2.4)
        add_rounded_rect(s, x, Inches(1.8), Inches(1.8), Inches(0.9),
                         color, label, font_size=13)
        if i < len(pipeline_stages) - 1:
            add_text_box(s, x + Inches(1.8), Inches(1.95), Inches(0.6), Inches(0.5),
                         "->", font_size=20, color=ACCENT, bold=True,
                         align=PP_ALIGN.CENTER)

    # 実行コマンド
    add_text_box(s, Inches(0.8), Inches(3.0), Inches(5), Inches(0.4),
                 "実行コマンド", font_size=16, color=ACCENT, bold=True)
    cmd_text = (
        "REPO_URL=...  BASE_REF=main\n"
        "PATCH_FILE=/tmp/drone-poc/patch.diff\n"
        "ARTIFACTS_ROOT=/tmp/drone-poc/artifacts\n"
        "python -m agent_runner.cli --goal \"kpを1.5に上げて応答性を改善\""
    )
    add_text_box(s, Inches(0.8), Inches(3.4), Inches(6.0), Inches(1.8),
                 cmd_text, font_size=11, color=ACCENT, font_name="Courier New")

    # 生成されたアーティファクト一覧
    add_text_box(s, Inches(7.2), Inches(3.0), Inches(5), Inches(0.4),
                 "生成 Artifacts", font_size=16, color=ACCENT, bold=True)
    artifact_items = [
        "request.json    - 目標・制約の記録",
        "git.json        - repo / base_ref / SHA",
        "patch.diff      - 適用パッチ",
        "patch_provider.json - パッチ提供元情報",
        "params.json     - seed / episodes / scene",
        "runtime.json    - host / Python版 / deps",
        "metrics.json    - baseline/candidate/delta",
        "summary.json    - 合否判定・比較結果",
        "evaluation_profile.json",
        "episodes_*.jsonl - raw observations",
    ]
    add_bullet_list(s, Inches(7.2), Inches(3.4), Inches(5.5), Inches(3.0),
                    artifact_items, font_size=11, color=LIGHT_GRAY)

    # 結果ステータスボックス
    add_rounded_rect(s, Inches(0.8), Inches(5.5), Inches(5.5), Inches(1.2),
                     RGBColor(0x2A, 0x4A, 0x2A),
                     "Status: SUCCEEDED\n"
                     "15 unit tests passed / lint passed\n"
                     "stub sim-eval でメトリクス生成・比較・合否判定まで完走",
                     font_size=13, font_color=GREEN_SOFT)

    # 注記
    add_text_box(s, Inches(7.2), Inches(5.8), Inches(5.5), Inches(0.8),
                 "* 現時点では eval/run.py はスタブ（ランダム値）\n"
                 "* CoppeliaSim 実機接続は次ステップで実装",
                 font_size=12, color=RGBColor(0xFF, 0xCC, 0x80))

    # ════════════════════════════════════════
    # Slide 8: Kubernetes案の検討と却下
    # ════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    slide_title(s, "Kubernetes案の検討と却下")

    # 左: K8s案
    add_text_box(s, Inches(0.8), Inches(1.5), Inches(5), Inches(0.5),
                 "当初の Kubernetes 案", font_size=20, color=ACCENT2, bold=True)
    k8s_items = [
        "Job による run 実行、CronJob による定期実行",
        "PVC / S3 による artifact 管理",
        "Helm chart によるデプロイ",
        "Prometheus / Grafana による監視",
        "複数 simulator の並列スケジューリング",
    ]
    add_bullet_list(s, Inches(0.8), Inches(2.1), Inches(5.5), Inches(3),
                    k8s_items, font_size=14, color=LIGHT_GRAY)

    # 右: 却下理由
    add_text_box(s, Inches(7.0), Inches(1.5), Inches(5), Inches(0.5),
                 "却下理由", font_size=20, color=RED_SOFT, bold=True)
    reject_items = [
        "PoC段階で分散実行は不要",
        "CoppeliaSim は単一インスタンスで十分",
        "K8s のセットアップ・運用コストが高い",
        "まず「1 run が確実に完走する」ことが最優先",
        "再現性・比較可能性はローカルでも達成可能",
    ]
    add_bullet_list(s, Inches(7.0), Inches(2.1), Inches(5.5), Inches(3),
                    reject_items, font_size=14, color=LIGHT_GRAY)

    # 結論ボックス
    add_rounded_rect(s, Inches(2.5), Inches(5.2), Inches(8), Inches(1.2),
                     RGBColor(0x2A, 0x4A, 0x2A),
                     "結論: 単一ホスト + ローカルファイルで PoC を構築\n"
                     "将来の拡張時に K8s への移行パスは残す",
                     font_size=16, font_color=GREEN_SOFT)

    # ════════════════════════════════════════
    # Slide 7: 評価方法
    # ════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    slide_title(s, "評価方法")

    add_text_box(s, Inches(0.8), Inches(1.5), Inches(11), Inches(0.6),
                 "1 run = baseline と candidate の同一条件比較",
                 font_size=18, color=WHITE, bold=True)

    # Metrics table header
    add_rounded_rect(s, Inches(0.8), Inches(2.3), Inches(3.2), Inches(0.6),
                     ACCENT2, "メトリクス", font_size=14)
    add_rounded_rect(s, Inches(4.2), Inches(2.3), Inches(4.0), Inches(0.6),
                     ACCENT2, "定義", font_size=14)
    add_rounded_rect(s, Inches(8.4), Inches(2.3), Inches(2.5), Inches(0.6),
                     ACCENT2, "方向", font_size=14)

    metrics_data = [
        ("success_rate", "成功エピソード数 / 総エピソード数", "高いほど良い"),
        ("collision_count_mean", "全エピソードの衝突合計 / 総エピソード数", "低いほど良い"),
        ("time_to_goal_mean_sec", "成功エピソードのゴール到達平均時間", "低いほど良い"),
        ("reward_mean", "各エピソードの累積報酬の平均", "高い（診断用）"),
    ]
    for i, (name, defn, direction) in enumerate(metrics_data):
        y = Inches(3.0 + i * 0.55)
        bg = BG_MID if i % 2 == 0 else RGBColor(0x2A, 0x2A, 0x42)
        add_rounded_rect(s, Inches(0.8), y, Inches(3.2), Inches(0.5),
                         bg, name, font_size=12, font_color=ACCENT)
        add_rounded_rect(s, Inches(4.2), y, Inches(4.0), Inches(0.5),
                         bg, defn, font_size=11, font_color=LIGHT_GRAY)
        add_rounded_rect(s, Inches(8.4), y, Inches(2.5), Inches(0.5),
                         bg, direction, font_size=12, font_color=LIGHT_GRAY)

    # evaluation_profile
    add_text_box(s, Inches(0.8), Inches(5.5), Inches(11), Inches(1.2),
                 "evaluation_profile により評価判定を切り替え可能\n"
                 "- simulation 再実行なしに重み・閾値・主指標を変更できる\n"
                 "- raw observation (episodes JSONL) から再集計が可能",
                 font_size=14, color=LIGHT_GRAY)

    # ════════════════════════════════════════
    # Slide 8: ネクストアクション
    # ════════════════════════════════════════
    s = prs.slides.add_slide(blank)
    slide_title(s, "ネクストアクション")

    # 技術面
    add_text_box(s, Inches(0.8), Inches(1.5), Inches(5), Inches(0.5),
                 "技術面", font_size=20, color=ACCENT, bold=True)
    tech_items = [
        "CoppeliaSim 実機接続の検証",
        "eval/run.py を実際のシミュレーション評価に差し替え",
        "End-to-end で 1 run を完走させる",
    ]
    add_bullet_list(s, Inches(0.8), Inches(2.1), Inches(5.5), Inches(2.5),
                    tech_items, font_size=15, color=LIGHT_GRAY)

    # 研究面
    add_text_box(s, Inches(0.8), Inches(4.0), Inches(5), Inches(0.5),
                 "研究面", font_size=20, color=ORANGE, bold=True)
    research_items = [
        "関連論文の調査が不足している",
        "  - 他の研究ではどのような評価指標を用いているか",
        "  - LLM によるコード生成の評価手法の先行研究",
        "  - ドローン制御における sim-to-real gap の扱い",
        "論文を読む時間の確保が課題",
    ]
    add_bullet_list(s, Inches(0.8), Inches(4.6), Inches(11), Inches(2.5),
                    research_items, font_size=15, color=LIGHT_GRAY)

    # 課題ボックス
    add_rounded_rect(s, Inches(7.0), Inches(1.8), Inches(5.5), Inches(2.0),
                     RGBColor(0x5A, 0x2A, 0x2A),
                     "課題\n\n時間の確保が難しい状況\n研究・論文調査に充てる\nまとまった時間が必要",
                     font_size=14, font_color=RGBColor(0xFF, 0xCC, 0xCC))

    # ════════════════════════════════════════
    # Save
    # ════════════════════════════════════════
    output = "/Users/takagiyuuki/maeda_-laboratory/progress_report.pptx"
    prs.save(output)
    print(f"Generated: {output}")


if __name__ == "__main__":
    main()
